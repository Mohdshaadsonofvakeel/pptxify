# app/pptx_builder.py
from io import BytesIO
from typing import List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.dml.color import RGBColor

from .schemas import Outline
from .template_utils import extract_template_images, find_preferred_layout


def _set_title(slide, text: str):
    # Find a title placeholder and set text
    for shp in slide.placeholders:
        try:
            if shp.placeholder_format.type == PP_PLACEHOLDER.TITLE or shp.placeholder_format.type == PP_PLACEHOLDER.CENTER_TITLE:
                shp.text = text
                return
        except Exception:
            continue
    # Fallback: first placeholder
    if slide.placeholders:
        slide.placeholders[0].text = text


def _set_bullets(slide, bullets: list[str]):
    # Find a body/content placeholder
    body = None
    for shp in slide.placeholders:
        try:
            if shp.placeholder_format.type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.CONTENT, PP_PLACEHOLDER.SUBTITLE):
                body = shp
                break
        except Exception:
            continue

    if body is None:
        # nothing to do
        return

    tf = body.text_frame
    tf.clear()
    if not bullets:
        return

    # First bullet
    p = tf.paragraphs[0]
    p.level = 0
    p.text = bullets[0]

    # Others
    for b in bullets[1:]:
        r = tf.add_paragraph()
        r.text = b
        r.level = 0


def _first_picture_placeholder(slide):
    for shp in slide.placeholders:
        try:
            if shp.placeholder_format.type in (PP_PLACEHOLDER.PICTURE, PP_PLACEHOLDER.CONTENT):
                # This placeholder supports inserting a picture
                return shp
        except Exception:
            continue
    return None


def _insert_picture(slide, image_bytes: bytes):
    ph = _first_picture_placeholder(slide)
    if ph is not None:
        try:
            from io import BytesIO
            ph.insert_picture(BytesIO(image_bytes))
            return True
        except Exception:
            pass
    # Fallback: add picture on the right side if no placeholder
    try:
        from io import BytesIO
        slide.shapes.add_picture(BytesIO(image_bytes), Inches(6.0), Inches(1.5), width=Inches(3.0))
        return True
    except Exception:
        return False


def build_presentation(outline: Outline, template_bytes: bytes) -> bytes:
    prs = Presentation(BytesIO(template_bytes))
    template_images = extract_template_images(template_bytes)
    img_idx = 0

    # Title slide if layout exists
    title_layout = find_preferred_layout(prs, ["Title Slide", "Title Only", "Section Header", "Title"])
    if title_layout is not None:
        title_slide = prs.slides.add_slide(title_layout)
        _set_title(title_slide, outline.title)
    else:
        # Fall back to first layout
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        _set_title(slide, outline.title)

    # Content slides
    for s in outline.slides:
        layout_name = s.layout.lower() if s.layout else "auto"
        if layout_name == "auto":
            layout = find_preferred_layout(prs, ["Title and Content", "Content with Caption", "Two Content", "Picture with Caption", "Blank"])
        else:
            layout = find_preferred_layout(prs, [s.layout, "Title and Content", "Content with Caption", "Two Content"])

        if layout is None:
            layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]

        slide = prs.slides.add_slide(layout)
        _set_title(slide, s.title)
        _set_bullets(slide, s.bullets)

        # Reuse template images opportunistically
        if template_images:
            inserted = _insert_picture(slide, template_images[img_idx % len(template_images)])
            if inserted:
                img_idx += 1

        # Speaker notes (optional)
        if s.notes is not None:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.clear()
            notes_slide.notes_text_frame.text = s.notes or ""

    bio = BytesIO()
    prs.save(bio)
    return bio.getvalue()
